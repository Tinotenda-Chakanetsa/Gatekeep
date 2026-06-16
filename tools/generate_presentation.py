#!/usr/bin/env python3
"""Generate the Gatekeeper oral-presentation slide deck (.pptx).

Designed for a 10-15 minute oral defence scored against this rubric:
    Presentation manner            3
    Arrangement of presentation    6
    Accuracy and clarity           6
    Understanding of project       9    (largest weight)
    Confidence in delivery         3
    Managing question and answers  3
    --------------------------------
    TOTAL                         30

Every slide is plain text-frames + speaker notes — no embedded images or fancy
SmartArt — so the deck remains fully editable in PowerPoint without losing
formatting on save.

Run with the Python that has python-pptx installed:
    python tools/generate_presentation.py
Output: Gatekeeper_Presentation.pptx in the project root.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PATH = PROJECT_ROOT / "Gatekeeper_Presentation.pptx"

# Real training-run artefacts that we embed directly.
TRAINING_RUN = Path(
    r"C:\Users\Admin\Desktop\OCR project\Lisence plate detector prototype"
    r"\license_plate_trainer_package\trained\runs\license_plate_training"
)

ACCENT = RGBColor(0x2E, 0x6B, 0x4F)        # dashboard green
INK = RGBColor(0x1B, 0x1B, 0x1B)
MUTED = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF3, 0xF0, 0xE6)      # matches dashboard cream
PLACEHOLDER_FILL = RGBColor(0xEE, 0xEA, 0xDD)
PLACEHOLDER_BORDER = RGBColor(0xB0, 0xA8, 0x90)


# ---------------------------------------------------------------------------
# slide helpers
# ---------------------------------------------------------------------------

def _add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _style_run(run, *, size, bold=False, color=INK, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color


def add_title_bar(slide, title: str, kicker: str | None = None):
    # accent stripe down the left
    stripe = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    tf = _add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12.5), Inches(1.1))
    if kicker:
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = kicker.upper()
        _style_run(r, size=11, bold=True, color=MUTED)
        p.space_after = Pt(2)
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    r = p2.add_run()
    r.text = title
    _style_run(r, size=30, bold=True, color=INK)


def add_bullets(slide, items, *, left=0.65, top=1.65, width=12.0, height=5.4,
                size=18, bullet_color=ACCENT):
    tf = _add_textbox(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            heading_text, sub_text = item
        else:
            heading_text, sub_text = item, None

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6 if i > 0 else 0)
        p.space_after = Pt(2)

        bullet = p.add_run()
        bullet.text = "• "
        _style_run(bullet, size=size, bold=True, color=bullet_color)

        run = p.add_run()
        run.text = heading_text
        _style_run(run, size=size, bold=True, color=INK)

        if sub_text:
            sub = tf.add_paragraph()
            sub.space_after = Pt(4)
            r = sub.add_run()
            r.text = "      " + sub_text
            _style_run(r, size=size - 4, color=MUTED, italic=True)


def add_speaker_notes(slide, text: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text.strip()


def add_image(slide, image_path: Path, *, left, top, width=None, height=None):
    """Embed a real image. Falls back to a placeholder if the file is missing."""
    if not image_path.exists():
        add_image_placeholder(slide, left=left, top=top,
                              width=width or 6, height=height or 4,
                              caption=f"[ MISSING: {image_path.name} ]")
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)


def add_image_placeholder(slide, *, left, top, width, height, caption: str):
    """Draw a clearly-marked drop-zone for an image the user will paste in
    later via PowerPoint's Insert > Picture, or by right-clicking the shape
    and choosing Change Picture."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PLACEHOLDER_FILL
    box.line.color.rgb = PLACEHOLDER_BORDER
    box.line.width = Pt(1.5)
    box.line.dash_style = 7  # dashed

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[ INSERT IMAGE HERE ]"
    _style_run(r, size=14, bold=True, italic=True, color=ACCENT)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = caption
    _style_run(r2, size=11, italic=True, color=MUTED)


# ---------------------------------------------------------------------------
# the deck
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # fully blank — we draw our own

    # ---------- 1. Title ----------
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    accent = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    tf = _add_textbox(s, Inches(1.1), Inches(1.5), Inches(11.5), Inches(4.5))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "GATEKEEPER"
    _style_run(r, size=58, bold=True, color=INK)

    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "An Automated Licence-Plate Gate Access-Control System\nUsing Deep Learning and IoT"
    _style_run(r, size=24, color=ACCENT)

    p = tf.add_paragraph()
    p.space_before = Pt(28)
    r = p.add_run()
    r.text = "Industrial Training Oral Presentation  ·  June 2026"
    _style_run(r, size=14, italic=True, color=MUTED)

    info_tf = _add_textbox(s, Inches(1.1), Inches(5.4), Inches(11.5), Inches(1.8))
    info_lines = [
        ("Student:", "Tinotenda M. Chakanetsa"),
        ("Registration No.:", "R235191W"),
        ("Programme:", "BSc Honours Cloud Computing and Internet of Things"),
        ("Faculty:", "Computer Engineering, Informatics and Communications"),
        ("Supervisor:", "Mr R. Zenda"),
    ]
    for i, (label, value) in enumerate(info_lines):
        if i == 0:
            p = info_tf.paragraphs[0]
        else:
            p = info_tf.add_paragraph()
        p.space_after = Pt(1)
        rl = p.add_run()
        rl.text = label + " "
        _style_run(rl, size=12, bold=True, color=INK)
        rv = p.add_run()
        rv.text = value
        _style_run(rv, size=12, color=INK)

    add_speaker_notes(s, """
Good morning, panel. My name is Tinotenda Chakanetsa, registration number
R235191W, third-year BSc Honours Cloud Computing and IoT. My internship
project is Gatekeeper — an automated licence-plate gate access-control
system built around deep learning and a small IoT bill of materials.
I'll spend roughly twelve minutes walking you through the problem,
how the system is built, what was actually achieved, and where I would
take it next. Please hold questions for the end — I'll be glad to take
them when we get to the Q&A slide.
""")

    # ---------- 2. Agenda ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Agenda", kicker="What I will cover")
    add_bullets(s, [
        ("Problem and objectives", "Why this project, and what it set out to deliver"),
        ("System architecture and hardware build", "ESP32 + ESP32-CAM + servo + LEDs + buzzer"),
        ("Dataset, model training and training visuals", "Roboflow merge + YOLO11n fine-tune"),
        ("Backend, frontend and firmware implementation", "How each layer is built"),
        ("OCR live capture and dashboard", "A real plate, end-to-end"),
        ("Deployment, results, limitations, next steps", "Free-tier stack, what works, what is rough"),
        ("Questions", "Open discussion"),
    ], size=18)
    add_speaker_notes(s, """
This is the route through the talk. I'll spend more time on architecture
and results — those are where the engineering decisions and the
measurable outcomes live — and keep background, deployment, and future
work tight. Total budget is around fifteen minutes including questions.
""")

    # ---------- 3. Background & Problem ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Background and Problem Statement", kicker="Chapter 1")
    add_bullets(s, [
        ("Manual gate access dominates in Zimbabwe.",
         "Guards verify vehicles by eye and record plates by hand — if at all."),
        ("It is slow, inconsistent and leaves no real audit trail.",
         "After-the-fact investigations rely on memory or unreliable CCTV."),
        ("Commercial ALPR exists but is priced for fleets and enterprises.",
         "Many require dedicated PC hardware at the gate and closed vendor clouds."),
        ("There is a gap for an open, low-cost, locally-relevant prototype.",
         "Something a single operator could deploy on commodity infrastructure."),
    ], size=20)
    add_speaker_notes(s, """
The starting observation is mundane but consequential. Manual gate
control is the default at almost every Zimbabwean residential complex,
office park, school and logistics yard I have seen. It is slow,
inconsistent and offers no searchable record. Commercial automated
licence-plate recognition exists, but it is sold to fleets and large
estates — not to a single homeowner or a small business. So there's a
clean gap to fill: a low-cost, open-source, locally-relevant prototype
that one operator can stand up themselves.
""")

    # ---------- 4. Aim & Objectives ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Aim and Objectives", kicker="Section 1.3")

    aim_tf = _add_textbox(s, Inches(0.65), Inches(1.55), Inches(12), Inches(1.0))
    p = aim_tf.paragraphs[0]
    r = p.add_run()
    r.text = "Aim:  "
    _style_run(r, size=18, bold=True, color=ACCENT)
    r = p.add_run()
    r.text = ("Design, build and evaluate a working low-cost end-to-end "
              "prototype of an automated licence-plate gate access-control system "
              "using deep-learning vision and IoT hardware.")
    _style_run(r, size=18, color=INK)

    add_bullets(s, [
        "Curate a Zimbabwean-context plate dataset by merging open Roboflow sources.",
        "Train and evaluate a YOLO11n detector that locates plates reliably.",
        "Implement a secure REST backend (Django + DRF + JWT) with role-based admin.",
        "Build a React dashboard for People, Vehicles, Access Logs and manual override.",
        "Write ESP32-CAM + ESP32 firmware that captures, decides and actuates the gate.",
        "Deploy the whole stack on free / commodity infrastructure.",
    ], top=2.85, size=17)
    add_speaker_notes(s, """
The single sentence at the top is the project aim. Below it are the six
objectives, in the order they were tackled. Each one became a chapter
of the implementation: dataset, model, backend, frontend, firmware,
deployment. I'll return to each of these in turn over the next few
slides.
""")

    # ---------- 5. System Architecture ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "System Architecture", kicker="The end-to-end picture")

    diagram = _add_textbox(s, Inches(0.65), Inches(1.5), Inches(12), Inches(3.6))
    code = ("[ PIR ]  ->  [ ESP32 control node ]  --UART CAPTURE-->  [ ESP32-CAM ]\n"
            "                       ^                                       |  WiFi/HTTPS\n"
            "                       | AUTH:1 / AUTH:0 / OPEN                 v\n"
            "                       +--UART---------------------+   [ Tailscale Funnel ]\n"
            "                                                              |\n"
            "                                                              v\n"
            "                                                       [ Django backend ]\n"
            "                                                       YOLO + PaddleOCR\n"
            "                                                       SQLite + media volume\n"
            "                                                              ^\n"
            "[ Browser ] -- HTTPS --> [ Vercel React SPA ] -- CORS --> [ Tailscale Funnel ]")
    p = diagram.paragraphs[0]
    r = p.add_run()
    r.text = code
    _style_run(r, size=12, color=INK, font="Consolas")

    add_bullets(s, [
        ("Three physical zones.", "Gate site, public HTTPS hop, backend VM."),
        ("Two-board gate node.", "Control node owns peripherals; camera node owns WiFi/HTTPS."),
        ("Dashboard talks to backend directly.", "Browser ↔ Vercel ↔ Funnel ↔ Django, CORS-enabled."),
    ], top=5.0, size=16)
    add_speaker_notes(s, """
This is the single slide I am happy to dwell on if asked. There are
three physical zones: the gate site at the top, the public-internet hop
in the middle, and the backend host at the bottom. The two boards at the
gate cooperate over a UART link — the ESP32 control node owns the PIR,
the servo boom, the LEDs and the buzzer; the ESP32-CAM owns WiFi, HTTPS
and the camera. The administrative path runs in parallel: a browser
hits the React SPA on Vercel, Vercel hits the backend through the
Tailscale Funnel. The funnel exists because the VM has no public IPv4 of
its own.
""")

    # ---------- 6. Hardware Build ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Hardware Build", kicker="The gate node, in person")
    add_image_placeholder(
        s, left=0.45, top=1.65, width=7.0, height=5.3,
        caption="Prototype photo — ESP32 + ESP32-CAM + PIR + SG90 servo + green/red LEDs + buzzer wired on a breadboard, with the shared 5 V supply visible.",
    )
    add_bullets(
        s,
        [
            ("ESP32 control node.", "PIR D13, servo D25, green D26, red D27, buzzer D33."),
            ("ESP32-CAM camera node.", "OV2640 capture + WiFi/HTTPS to the backend."),
            ("UART link, 3 wires + GND.", "CAM IO14 → D16(RX2); CAM IO15 ← D17(TX2); common ground."),
            ("Bill of materials ≈ USD 20.", "Inside the original low-cost objective."),
            ("Power: shared 5 V 2 A rail.", "Servo + buzzer + both boards; never drive servo from 3V3."),
        ],
        left=7.7, top=1.65, width=5.4, height=5.3, size=14,
    )
    add_speaker_notes(s, """
This is the physical build. Two boards working together at the gate
site. The ESP32 control board owns the peripherals you can see, the
ESP32-CAM owns the camera and the WiFi. They talk over a three-wire
UART link, plus a shared ground that is non-negotiable. Total cost of
all the parts on this breadboard is under twenty US dollars. Power is
delivered from a single 5 V 2 A supply that feeds both boards and the
servo — running the servo from the ESP32's 3.3 V pin will crash the
chip mid-capture, so the rail is shared but the regulator is bypassed.
""")

    # ---------- 7. Dataset ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Dataset", kicker="Section 3.3 + 4.1")
    add_bullets(s, [
        ("Two open Roboflow datasets, merged.",
         "ta-jfzoz/licence-detection-fy0ma v3   +   msucs/zimbabwe-license-plates v4"),
        ("All labels collapsed to a single class: license_plate.",
         "merge_detector_datasets.py walks each split and rewrites every label."),
        ("717 images total — 635 train / 55 valid / 27 test.",
         "Approximate 70 / 20 / 10 split, inherited from the Roboflow exports."),
        ("Single data.yaml + standard YOLO layout.",
         "Ready to feed straight into Ultralytics train_detector.py."),
    ], size=20)
    add_speaker_notes(s, """
Rather than collecting plates from scratch — which would have eaten the
whole project budget — I merged two publicly available Roboflow
datasets. A small Python helper downloads each, normalises the file
layout, and rewrites every label so that whatever class IDs the upstream
used, the merged dataset has exactly one class: license_plate. After
merging, I had 717 images: 635 training, 55 validation, 27 test. That
follows roughly the 70 / 20 / 10 convention.
""")

    # ---------- 7. Model Training ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Model Training", kicker="Section 4.2 + 5.1")
    add_bullets(s, [
        ("Base: YOLO11n, transfer-learned from yolo11n.pt (COCO).",
         "Smallest Ultralytics backbone — trains in minutes on a single GPU."),
        ("100 max epochs, patience = 25, image size 640, batch 16, AMP on.",
         "Effective 75 epochs after early-stop; ~6.5 minutes wall clock."),
        ("Best validation: mAP@0.5 ≈ 0.94, mAP@0.5:0.95 ≈ 0.65.",
         "Precision peaked at 1.00; recall around 0.90."),
        ("Best checkpoint deployed as artifacts/license_plate_detector.pt.",
         "Loaded lazily by the Django backend on first capture."),
    ], size=20)
    add_speaker_notes(s, """
I picked the smallest Ultralytics backbone, YOLO11n, because the dataset
is modest in size and the inference will run on a CPU-only VM anyway.
Transfer learning from the COCO checkpoint accelerated convergence and
gave a much stronger starting point than random initialisation.
Training ran for an effective 75 epochs in about six and a half minutes
on a single GPU. The headline number is validation mAP at 0.5
intersection-over-union of approximately 0.94 — well into the "usable
for the next OCR stage" range. The mAP@0.5:0.95 figure of around 0.65
tells me the boxes are well-localised but not razor-tight, which is
fine because I pad the crop before passing it to PaddleOCR anyway.
""")

    # ---------- 9. Training Visuals (real images on disk) ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Training Visuals", kicker="What the trainer reported")
    # Left: loss + metric curves over the 75 effective epochs
    add_image(s, TRAINING_RUN / "results.png",
              left=0.45, top=1.6, width=7.6, height=4.0)
    # Right: predictions on a validation batch (qualitative ground truth)
    add_image(s, TRAINING_RUN / "val_batch0_pred.jpg",
              left=8.3, top=1.6, width=4.5, height=4.0)
    # Bottom strip: small confusion matrix + a short caption
    add_image(s, TRAINING_RUN / "confusion_matrix_normalized.png",
              left=0.45, top=5.85, width=3.0, height=1.5)
    caption = _add_textbox(s, Inches(3.75), Inches(5.9), Inches(9.1), Inches(1.4))
    p = caption.paragraphs[0]
    r = p.add_run()
    r.text = ("Left: per-epoch losses and validation metrics — losses fall monotonically; "
              "mAP@0.5 plateaus around 0.94 in the second half of training.   "
              "Right: detector boxes drawn on real validation imagery.   "
              "Below-left: normalised confusion matrix (single class).")
    _style_run(r, size=12, italic=True, color=MUTED)
    add_speaker_notes(s, """
This is the slide for the panel that wants to see the model justified.
The graph on the left shows training and validation losses against the
two mAP curves over the 75 effective epochs — the losses fall
monotonically and the metrics plateau, exactly the early-stop signal we
hoped for. The picture on the right shows the detector boxes drawn on
real validation imagery so the panel can see the model finding plates
in cluttered scenes, not just numbers. The small panel at the bottom
left is the normalised confusion matrix — single class, so the only
off-diagonal cell is background false positives, which is small.
""")

    # ---------- 10. Backend ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Backend Implementation", kicker="Section 4.3")
    add_bullets(s, [
        ("Django 5 + DRF + djangorestframework-simplejwt.",
         "Three apps: accounts (custom User + role), gate (domain + API), ocr_api (vision)."),
        ("Domain models: Person, Vehicle, GateDevice, AccessLog.",
         "Plus a custom User with admin / operator roles, enforced by DRF permissions."),
        ("Two authentication channels.",
         "JWT bearer token for dashboard users; X-Device-Key header for the ESP32-CAM."),
        ("Three-tier plate matching: exact → confusion-folded → bounded fuzzy.",
         "Tolerates common OCR noise (O/0, I/1, S/5, B/8) without false matches."),
        ("YOLO and PaddleOCR loaded lazily inside service functions.",
         "Keeps Django startup and migrations free of the multi-GB ML deps."),
    ], size=17)
    add_speaker_notes(s, """
The backend is a fairly conventional Django REST framework setup, split
into three apps. The interesting parts are two design choices. First,
two authentication channels: dashboard users get JWT tokens, while the
ESP32-CAM authenticates with a per-device API key in a header. That
keeps user passwords out of the firmware entirely. Second, plate
matching is layered: an exact match first, then a confusion-folded
match that collapses commonly-confused characters like letter-O and
zero, and finally a bounded fuzzy match. The OCR is the noisy stage of
the pipeline, so the matching has to absorb that noise without granting
the wrong vehicle.
""")

    # ---------- 11. Frontend ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Frontend Dashboard", kicker="Section 4.4")
    add_bullets(
        s,
        [
            ("React 19 + Vite + react-router 7, on Vercel.",
             "Static SPA served from Vercel's edge."),
            ("Role-aware nav and protected routes.",
             "Operator: browse + override.   Admin: mutate registry."),
            ("Pages.",
             "Dashboard, Vehicles, People, Access Logs, Devices, Users, OCR Tester."),
            ("VITE_API_BASE_URL + CORS.",
             "Vercel reserves /api; SPA points straight at the funnel hostname."),
        ],
        left=0.55, top=1.65, width=5.6, height=5.4, size=14,
    )
    add_image_placeholder(
        s, left=6.35, top=1.65, width=6.6, height=5.4,
        caption="Dashboard screenshot — the Gate Control Dashboard page showing the four metric cards (Registered Vehicles, Granted Today, Denied Today, Devices Online) plus the live Recent Gate Events table with capture thumbnails.",
    )
    add_speaker_notes(s, """
The frontend is a React 19 single-page app, hosted statically on
Vercel's Hobby tier. The interesting story here is the deployment
quirk I hit: Vercel reserves any URL path containing the segment slash
api for its own serverless functions. My first deployment tried to
proxy /api back to the backend through a Vercel rewrite and that hit a
404 wall. The fix was to point the SPA at the backend's absolute URL
through a build-time environment variable and enable CORS for the
Vercel origin on the Django side. Worth knowing if anyone else hits it.
""")

    # ---------- 10. Firmware ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Firmware: Two-Board Split", kicker="Section 4.5")
    add_bullets(s, [
        ("ESP32 control node — peripherals only.",
         "Reads the PIR, drives SG90 servo boom, green/red LEDs, buzzer."),
        ("ESP32-CAM camera node — network + camera only.",
         "Captures JPEGs, POSTs to /api/gate/check/, parses the JSON decision."),
        ("Linked over UART at 115200 baud, line-terminated text.",
         "CAPTURE\\n    AUTH:1\\n    AUTH:0\\n    ERR\\n    OPEN\\n    READY\\n"),
        ("Solves the AI-Thinker GPIO scarcity problem cleanly.",
         "Camera board was running out of pins; this split puts each role on its own MCU."),
        ("Camera fixes: sensor un-mirror + two-frame buffer drain.",
         "set_hmirror(1) cancels the OV2640 default; drain ensures genuinely-fresh frames."),
    ], size=17)
    add_speaker_notes(s, """
The original prototype ran everything on one ESP32-CAM and quickly ran
out of GPIO once I started adding a servo, LEDs and a buzzer. So I
split the firmware across two boards that talk to each other over a
plain UART link at 115200 baud, exchanging short text lines like
CAPTURE, AUTH 1 and AUTH 0. The control node owns the peripherals and
the camera node owns the network. Two camera-specific gotchas I had to
solve along the way are worth flagging: the AI-Thinker board's default
sensor output is mirrored, so I have to call set_hmirror to un-mirror
it before PaddleOCR sees the frame; and the camera buffers were
producing stale frames, so I drain two stale frames before grabbing
the real one. Once those two fixes were in, end-to-end accuracy was
solid.
""")

    # ---------- 13. Deployment ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Deployment", kicker="Section 4.6")
    add_bullets(
        s,
        [
            ("Backend in Docker on an Ubuntu VM (private IPv4 only).",
             "docker compose; DB + media in named volumes."),
            ("Tailscale Funnel: public HTTPS to localhost:8000.",
             "Free, no public IP, Let's Encrypt cert on *.ts.net."),
            ("Vercel auto-deploys the SPA on push to main.",
             "GitHub integration; zero-ops releases."),
            ("Recurring cost: just the VM.",
             "Vercel Hobby + Tailscale personal are free."),
        ],
        left=0.55, top=1.65, width=5.6, height=5.4, size=14,
    )
    add_image_placeholder(
        s, left=6.35, top=1.65, width=6.6, height=2.55,
        caption="Vercel Deployments tab — current Production build tied to the GitHub commit (gatekeep-kappa.vercel.app).",
    )
    add_image_placeholder(
        s, left=6.35, top=4.45, width=6.6, height=2.55,
        caption="Tailscale admin / `tailscale funnel status` output — VM enrolled with Funnel enabled at https://user-hvm-domu.taile2f0da.ts.net.",
    )
    add_speaker_notes(s, """
Deployment is where I deliberately leaned on free tiers. The backend
runs inside Docker on a private-IP Ubuntu VM. Because the VM has no
public IPv4, I use Tailscale Funnel to expose port 8000 over a public
HTTPS URL — no router config, no DNS, no certbot — just one cloudflared-
style daemon that gives me a permanent .ts.net hostname with a real
Let's Encrypt cert. The React frontend is deployed on Vercel Hobby,
which auto-deploys on every push to main. The whole running cost is
just the VM itself.
""")

    # ---------- 12. Results ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Results", kicker="Chapter 5")
    add_bullets(s, [
        ("Detector: mAP@0.5 ≈ 0.94 on the held-out validation split.",
         "Peak precision 1.00, peak recall ≈ 0.90 at the best operating point."),
        ("Backend test suite: 21 of 21 checks pass.",
         "JWT login, role-based CRUD, plate normalisation, device-key auth, gate-command queue."),
        ("Capture-pipeline tests: 9 of 9 checks pass.",
         "Known plate → GRANTED; noisy variant resolves via fuzzy match; unknown → DENIED."),
        ("End-to-end live: PIR → image → OCR → match → boom rises.",
         "Confirmed in person with a real Honda plate."),
        ("Bill of materials: ~USD 20 per gate plus the VM.",
         "Within the original \"low-cost\" objective."),
    ], size=18)
    add_speaker_notes(s, """
Headline numbers: detector validation mAP at 0.5 is approximately 0.94,
the backend integration test suite passes all 21 of its checks
covering authentication, role-based access, plate normalisation and the
gate-command queue, and a separate capture-pipeline test that stubs the
OCR call passes all 9 of its assertions. The end-to-end system runs in
person — the live demo earlier in the term showed PIR triggering, the
cam uploading the plate, the backend granting, and the boom rising.
The total bill of materials at the gate is under twenty US dollars,
comfortably inside the low-cost objective.
""")

    # ---------- 15. Live Capture Example ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Live Capture Example", kicker="One real plate, end-to-end")
    add_image_placeholder(
        s, left=0.45, top=1.65, width=7.5, height=5.4,
        caption="OCR Tester result — original photo, YOLO-annotated detection box, cropped plate region, and recognised text returned by PaddleOCR. Captured from the dashboard's OCR Tester page.",
    )
    add_bullets(
        s,
        [
            ("Input.",
             "Frame captured by the ESP32-CAM on a PIR trigger."),
            ("YOLO11n.",
             "Locates the plate, returns a confidence-scored bounding box."),
            ("Crop + preprocess.",
             "Pad, optional grayscale + invert + 2× upscale before OCR."),
            ("PaddleOCR.",
             "Reads characters; joined text is what the backend uses."),
            ("Plate matching.",
             "Three tiers — exact, confusion-folded, bounded fuzzy."),
            ("AccessLog row + decision returned.",
             "Image persisted, dashboard shows it within seconds."),
        ],
        left=8.2, top=1.65, width=4.8, height=5.4, size=13,
    )
    add_speaker_notes(s, """
This is the single most evidence-rich slide for "the system actually
works". The composite on the left is taken from the OCR Tester page on
the dashboard, which exposes the full pipeline as a diagnostic tool. It
shows the original photo, the YOLO bounding box on the plate region,
the cropped plate, and the recognised text — the same pipeline that
runs on every PIR-triggered capture, just made visible. The bullet
list on the right is the order of operations: capture, detect, crop,
OCR, match, log. That whole chain runs in under four seconds on the
CPU-only VM.
""")

    # ---------- 16. Limitations & Future Work ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Limitations and Future Work", kicker="Sections 5.4 + 6.2")
    add_bullets(s, [
        ("OCR is CPU-only — about 2–4 s per capture.",
         "Fine for gate use; not suitable for a tollway."),
        ("OV2640 sensor is weak in low light.",
         "Night use would need a white-LED or IR illuminator."),
        ("setInsecure() on the ESP32 TLS path is prototype only.",
         "Production deployment would pin the Tailscale root CA."),
        ("Next: an LCD readout at the gate showing the matched plate + owner.",
         "Plenty of free GPIO on the control node already."),
        ("Next: WebSocket access-log stream to remove the 4 s polling lag.",
         "Django Channels + EventSource on the SPA."),
    ], size=19)
    add_speaker_notes(s, """
I am honest about three real limitations. First, CPU-only OCR adds
about two to four seconds to each gate cycle — fine for cars stopping
at a boom but not for a toll. Second, the OV2640 sensor struggles in
low light and would need an extra illuminator at night. Third,
disabling certificate validation on the ESP32 is a prototype shortcut;
a production deployment would pin the Tailscale root CA. The two
concrete next steps I would take are an LCD at the gate showing the
matched plate and owner, and switching the dashboard from a four-second
poll to a WebSocket stream for snappier live events.
""")

    # ---------- 14. Conclusion + Q&A ----------
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Thank You — Questions?", kicker="Conclusion")

    summary_tf = _add_textbox(s, Inches(0.65), Inches(1.55), Inches(12), Inches(2.2))
    p = summary_tf.paragraphs[0]
    r = p.add_run()
    r.text = ("All six objectives were achieved: a curated 717-image dataset, a YOLO11n "
              "detector at ≈ 0.94 mAP@0.5, a JWT + role-based Django REST backend with "
              "three-tier plate matching, a Vercel-hosted React dashboard, a "
              "two-board ESP32 firmware split and a fully free-tier deployment.")
    _style_run(r, size=17, color=INK)

    # Anticipated Q&A panel
    qa_tf = _add_textbox(s, Inches(0.65), Inches(3.9), Inches(12), Inches(3.3))
    p = qa_tf.paragraphs[0]
    r = p.add_run()
    r.text = "Anticipated questions I am ready for"
    _style_run(r, size=14, bold=True, color=ACCENT)
    p.space_after = Pt(4)

    qa_items = [
        "Why YOLO11n instead of a bigger backbone?    → Dataset is small; nano avoids over-fitting and runs on CPU.",
        "How does plate matching tolerate OCR errors?  → Three-tier: exact → character-folded → bounded Levenshtein.",
        "Why Tailscale Funnel and not Cloudflare Tunnel? → No domain required; works with a private-IP VM out of the box.",
        "What stops an attacker forging a gate request? → Device API key per camera; rotateable from the Devices admin page.",
        "Why split the firmware across two boards?     → ESP32-CAM is GPIO-starved; split puts peripherals on a board that has room.",
    ]
    for line in qa_items:
        p = qa_tf.add_paragraph()
        p.space_after = Pt(2)
        bullet = p.add_run()
        bullet.text = "• "
        _style_run(bullet, size=13, bold=True, color=ACCENT)
        r = p.add_run()
        r.text = line
        _style_run(r, size=13, color=INK)

    add_speaker_notes(s, """
End on a confident note. Restate the six objectives in one sentence and
invite questions. The bullets on the slide are the questions I have
rehearsed quick, specific answers for — when one of them comes up, I
already have the response. For anything I cannot answer fully on the
spot, the rule is: state what I know, state what I would do to find
out, and offer to follow up. Do not bluff. Thank the panel for their
time.

CHEAT SHEET FOR Q&A:
- Latency budget: capture ~150ms + upload ~1.5s + OCR ~2.5s = ~4s total.
- Dataset breakdown: 635 train / 55 valid / 27 test = 717 total.
- Best metrics row: epoch 66, precision 1.00, recall 0.92, mAP50 0.940.
- Stack versions: Django 5.2, DRF 3.15, React 19, Vite 7, Ultralytics 8.3.
- ESP32-CAM camera: OV2640, SVGA 800x600, JPEG quality 12, PSRAM-resident.
- Cost: ~USD 20 BOM at the gate, VM is the only recurring cost.
""")

    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
